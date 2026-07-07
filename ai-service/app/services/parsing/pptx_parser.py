"""PPTX parser."""

from pptx import Presentation

from app.models.kb import DocumentPage


def _shape_text(shape) -> str:  # noqa: ANN001
    if not hasattr(shape, "text"):
        return ""
    return str(shape.text or "").strip()


async def parse_pptx(
    file_path: str, document_id: int, document_version_id: int
) -> list[DocumentPage]:
    """Extract one DocumentPage per slide."""
    prs = Presentation(file_path)
    pages: list[DocumentPage] = []

    for index, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = _shape_text(slide.shapes.title)

        texts: list[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text and text not in texts:
                texts.append(text)

        body = "\n".join(texts).strip()
        pages.append(
            DocumentPage(
                document_id=document_id,
                document_version_id=document_version_id,
                page_number=index,
                title=title or f"第 {index} 页",
                text=body,
                metadata_json={"slide_index": index},
            )
        )

    return pages
