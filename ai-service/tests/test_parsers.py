"""Document parser tests."""

from pathlib import Path

import pytest

from app.core.errors import AppError


@pytest.mark.asyncio
async def test_parse_pdf_extracts_text_pages(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    from app.services.parsing import pdf_parser

    class FakePage:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class FakeReader:
        is_encrypted = False
        pages = [
            FakePage("第一章\nPDF 内容"),
            FakePage("第二页内容"),
        ]

    monkeypatch.setattr(pdf_parser, "PdfReader", lambda _: FakeReader())
    file_path = tmp_path / "sample.pdf"
    file_path.write_bytes(b"%PDF")

    pages = await pdf_parser.parse_pdf(str(file_path), document_id=1, document_version_id=2)

    assert len(pages) == 2
    assert pages[0].document_id == 1
    assert pages[0].document_version_id == 2
    assert pages[0].page_number == 1
    assert pages[0].title == "第一章"
    assert pages[0].text == "第一章\nPDF 内容"


@pytest.mark.asyncio
async def test_parse_pdf_rejects_scanned_pdf(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    from app.services.parsing import pdf_parser

    class FakePage:
        def extract_text(self) -> str:
            return ""

    class FakeReader:
        is_encrypted = False
        pages = [FakePage()]

    monkeypatch.setattr(pdf_parser, "PdfReader", lambda _: FakeReader())
    file_path = tmp_path / "scanned.pdf"
    file_path.write_bytes(b"%PDF")

    with pytest.raises(AppError) as exc:
        await pdf_parser.parse_pdf(str(file_path), document_id=1, document_version_id=2)

    assert exc.value.code == "PARSE_UNSUPPORTED"


@pytest.mark.asyncio
async def test_parse_pptx_extracts_slide_text(tmp_path: Path) -> None:
    from pptx import Presentation

    from app.services.parsing.pptx_parser import parse_pptx

    file_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "课程导论"
    slide.placeholders[1].text = "机器学习基础"
    prs.save(str(file_path))

    pages = await parse_pptx(str(file_path), document_id=1, document_version_id=2)

    assert len(pages) == 1
    assert pages[0].title == "课程导论"
    assert "机器学习基础" in (pages[0].text or "")


@pytest.mark.asyncio
async def test_parse_docx_splits_by_heading_1(tmp_path: Path) -> None:
    from docx import Document

    from app.services.parsing.docx_parser import parse_docx

    file_path = tmp_path / "doc.docx"
    doc = Document()
    doc.add_heading("第一章", level=1)
    doc.add_paragraph("第一章内容")
    doc.add_heading("第二章", level=1)
    doc.add_paragraph("第二章内容")
    doc.save(str(file_path))

    pages = await parse_docx(str(file_path), document_id=1, document_version_id=2)

    assert [p.title for p in pages] == ["第一章", "第二章"]
    assert "第一章内容" in (pages[0].text or "")
    assert "第二章内容" in (pages[1].text or "")


@pytest.mark.asyncio
async def test_parse_xlsx_creates_page_per_sheet(tmp_path: Path) -> None:
    from openpyxl import Workbook

    from app.services.parsing.xlsx_parser import parse_xlsx

    file_path = tmp_path / "book.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "成绩"
    ws.append(["姓名", "分数"])
    ws.append(["张三", 95])
    second = wb.create_sheet("说明")
    second.append(["字段", "含义"])
    wb.save(file_path)

    pages = await parse_xlsx(str(file_path), document_id=1, document_version_id=2)

    assert [p.title for p in pages] == ["成绩", "说明"]
    assert "| 姓名 | 分数 |" in (pages[0].text or "")
    assert "| 张三 | 95 |" in (pages[0].text or "")


@pytest.mark.asyncio
async def test_parse_file_dispatches_by_extension(monkeypatch: pytest.MonkeyPatch) -> None:
    from app.models.kb import DocumentPage
    from app.services.parsing import dispatcher

    async def fake_pdf(file_path: str, document_id: int, document_version_id: int):
        return [
            DocumentPage(
                document_id=document_id,
                document_version_id=document_version_id,
                page_number=1,
                text=file_path,
            )
        ]

    monkeypatch.setattr(dispatcher, "parse_pdf", fake_pdf)

    pages = await dispatcher.parse_file(
        "x.pdf", ".PDF", document_id=10, document_version_id=20
    )

    assert len(pages) == 1
    assert pages[0].document_id == 10
    assert pages[0].document_version_id == 20
    assert pages[0].text == "x.pdf"


@pytest.mark.asyncio
async def test_parse_file_rejects_unsupported_extension() -> None:
    from app.services.parsing.dispatcher import parse_file

    with pytest.raises(AppError) as exc:
        await parse_file("x.zip", ".zip", document_id=1, document_version_id=2)

    assert exc.value.code == "PARSE_UNSUPPORTED_FORMAT"


@pytest.mark.asyncio
async def test_parse_retry_success_after_retries() -> None:
    from app.services.parsing.retry import parse_retry

    attempts = 0

    @parse_retry(retries=2, base_delay=0, timeout=1)
    async def flaky() -> str:
        nonlocal attempts
        attempts += 1
        if attempts < 3:
            raise ConnectionError("temporary")
        return "ok"

    assert await flaky() == "ok"
    assert attempts == 3


@pytest.mark.asyncio
async def test_parse_retry_gives_up() -> None:
    from app.services.parsing.retry import parse_retry

    attempts = 0

    @parse_retry(retries=2, base_delay=0, timeout=1)
    async def always_fails() -> None:
        nonlocal attempts
        attempts += 1
        raise ConnectionError("still down")

    with pytest.raises(ConnectionError):
        await always_fails()

    assert attempts == 3


@pytest.mark.asyncio
async def test_parse_retry_does_not_retry_app_error() -> None:
    from app.services.parsing.retry import parse_retry

    attempts = 0

    @parse_retry(retries=2, base_delay=0, timeout=1)
    async def domain_error() -> None:
        nonlocal attempts
        attempts += 1
        raise AppError("PARSE_UNSUPPORTED", "unsupported")

    with pytest.raises(AppError):
        await domain_error()

    assert attempts == 1
